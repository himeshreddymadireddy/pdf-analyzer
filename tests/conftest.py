"""Shared deterministic builders for document-pipeline tests.

The builders return bytes and do not require a network or provider credentials.  Tests
that exercise actual PDF/PPTX parsing may import their optional dependencies locally.
"""
from __future__ import annotations

from io import BytesIO


def make_pdf_bytes(page_texts: list[str]) -> bytes:
    """Build a tiny invented-content PDF when PyMuPDF is installed."""
    try:
        import pymupdf
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required to build a PDF fixture.") from exc
    document = pymupdf.open()
    try:
        for text in page_texts:
            page = document.new_page()
            page.insert_text((72, 72), text)
        return document.tobytes(garbage=4, deflate=True)
    finally:
        document.close()


def make_pptx_bytes(slide_texts: list[str]) -> bytes:
    """Build a presentation with deterministic text-only slides."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required to build a PPTX fixture.") from exc
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    for text in slide_texts:
        slide = presentation.slides.add_slide(blank_layout)
        text_box = slide.shapes.add_textbox(0, 0, 5_000_000, 1_000_000)
        text_box.text_frame.text = text
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()
